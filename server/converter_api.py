# converter_api.py
#
# QuadraConverter - production-oriented conversion API
#
# Supports:
#   - Office -> PDF
#   - HTML -> PDF
#   - PDF -> Word
#   - PDF -> Excel
#   - PDF -> PowerPoint
#   - PDF -> PDF/A
#   - PDF Unlock
#   - PDF Protect
#   - PDF Translate
#   - Send converted file by email
#
# PDF -> Word:
#   Native PDF text extraction
#   Coordinate-aware line reconstruction
#   Font-size estimation
#   Scanned-PDF OCR fallback
#
# PDF -> Excel:
#   Native table detection
#   Line/coordinate extraction
#   Column clustering
#   Scanned-PDF OCR fallback
#
# IMPORTANT:
# This file intentionally contains ONLY ONE /convert route.
#
# Required Python packages:
#   fastapi
#   uvicorn
#   python-multipart
#   PyMuPDF
#   pdfplumber
#   python-docx
#   openpyxl
#   python-pptx
#   Pillow
#   pytesseract
#   resend
#
# Required system programs:
#   LibreOffice
#   qpdf
#   Ghostscript
#   Tesseract OCR
#
# ============================================================


import os
import re
import json
import shutil
import subprocess
import tempfile
import urllib.request
from pathlib import Path
from io import BytesIO
from typing import Any
from urllib.parse import quote

from fastapi import (
    FastAPI,
    File,
    Form,
    HTTPException,
    UploadFile,
)

from fastapi.middleware.cors import CORSMiddleware

from fastapi.responses import (
    FileResponse,
    JSONResponse,
)

from starlette.background import BackgroundTask


# ============================================================
# CONFIGURATION
# ============================================================

APP_NAME = "QuadraConverter Conversion API"
APP_VERSION = "3.0.0"


MAX_FILE_BYTES = int(
    os.getenv(
        "MAX_FILE_BYTES",
        str(100 * 1024 * 1024),
    )
)


OCR_DPI = int(
    os.getenv(
        "OCR_DPI",
        "250",
    )
)


OCR_MIN_TEXT_CHARS = int(
    os.getenv(
        "OCR_MIN_TEXT_CHARS",
        "40",
    )
)


CONVERSION_TIMEOUT = int(
    os.getenv(
        "CONVERSION_TIMEOUT",
        "600",
    )
)


MAX_PDF_PAGES = int(
    os.getenv(
        "MAX_PDF_PAGES",
        "500",
    )
)


CORS_ORIGINS = [
    origin.strip()
    for origin in os.getenv(
        "CORS_ORIGINS",
        "*",
    ).split(",")
    if origin.strip()
]


ALLOW_CREDENTIALS = (
    "*" not in CORS_ORIGINS
)


# ============================================================
# SUPPORTED FILE TYPES
# ============================================================

ALLOWED_PDF = {
    ".pdf",
}


ALLOWED_OFFICE = {
    ".doc",
    ".docx",
    ".ppt",
    ".pptx",
    ".xls",
    ".xlsx",
    ".odt",
    ".ods",
    ".odp",
    ".rtf",
}


ALLOWED_HTML = {
    ".html",
    ".htm",
    ".xhtml",
}


# ============================================================
# FASTAPI APP
# ============================================================

app = FastAPI(
    title=APP_NAME,
    version=APP_VERSION,
)


app.add_middleware(
    CORSMiddleware,
    allow_origins=CORS_ORIGINS,
    allow_credentials=ALLOW_CREDENTIALS,
    allow_methods=[
        "GET",
        "POST",
        "OPTIONS",
    ],
    allow_headers=["*"],
    expose_headers=[
        "Content-Disposition",
        "X-Converted-Filename",
        "X-Conversion-Engine",
    ],
)


# ============================================================
# GENERAL HELPERS
# ============================================================


def cleanup(path: Path | None):
    if path is None:
        return

    try:
        shutil.rmtree(
            path,
            ignore_errors=True,
        )
    except Exception:
        pass


def binary_path(
    *names: str,
) -> str:

    for name in names:

        found = shutil.which(name)

        if found:
            return found

    raise RuntimeError(
        "Required converter is not installed: "
        + ", ".join(names)
    )


def command_exists(
    *names: str,
) -> bool:

    for name in names:

        if shutil.which(name):
            return True

    return False


def safe_filename(
    filename: str | None,
    fallback: str = "converted-file",
) -> str:

    name = (
        Path(filename or fallback).name
    )

    name = re.sub(
        r"[^A-Za-z0-9._() \-]+",
        "_",
        name,
    )

    name = name.strip(
        " ."
    )

    return name or fallback


def validate_file_size(
    path: Path,
):

    if not path.exists():
        raise HTTPException(
            status_code=422,
            detail="Uploaded file was not saved.",
        )

    size = path.stat().st_size

    if size <= 0:
        raise HTTPException(
            status_code=400,
            detail="The uploaded file is empty.",
        )

    if size > MAX_FILE_BYTES:
        raise HTTPException(
            status_code=413,
            detail=(
                f"File exceeds the "
                f"{MAX_FILE_BYTES // (1024 * 1024)} MB "
                f"conversion limit."
            ),
        )


def save_upload(
    upload: UploadFile,
    work: Path,
    allowed: set[str],
) -> Path:

    original_name = (
        upload.filename
        or "input"
    )

    suffix = (
        Path(original_name)
        .suffix
        .lower()
    )

    if suffix not in allowed:

        raise HTTPException(
            status_code=400,
            detail=(
                f"Unsupported file type: "
                f"{suffix or 'unknown'}"
            ),
        )

    filename = safe_filename(
        original_name,
        f"input{suffix}",
    )

    source = (
        work / filename
    )

    size = 0

    try:

        with source.open(
            "wb"
        ) as output:

            while True:

                chunk = upload.file.read(
                    1024 * 1024
                )

                if not chunk:
                    break

                size += len(chunk)

                if size > MAX_FILE_BYTES:

                    raise HTTPException(
                        status_code=413,
                        detail=(
                            f"File exceeds the "
                            f"{MAX_FILE_BYTES // (1024 * 1024)} MB "
                            f"conversion limit."
                        ),
                    )

                output.write(
                    chunk
                )

    except HTTPException:
        raise

    except Exception as exc:

        raise HTTPException(
            status_code=400,
            detail=(
                f"Could not save uploaded file: "
                f"{exc}"
            ),
        )

    validate_file_size(
        source
    )

    return source


def run_checked(
    args: list[str],
    timeout: int = CONVERSION_TIMEOUT,
):

    try:

        process = subprocess.run(
            args,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            timeout=timeout,
            check=False,
        )

    except FileNotFoundError as exc:

        raise HTTPException(
            status_code=503,
            detail=(
                "Required conversion engine "
                f"is not installed: {exc}"
            ),
        )

    except subprocess.TimeoutExpired:

        raise HTTPException(
            status_code=504,
            detail=(
                "Conversion timed out. "
                "Try a smaller or simpler file."
            ),
        )

    stdout = (
        process.stdout or ""
    ).strip()

    stderr = (
        process.stderr or ""
    ).strip()

    if process.returncode != 0:

        detail = (
            stderr
            or stdout
            or "Conversion engine failed."
        )

        raise HTTPException(
            status_code=422,
            detail=detail[-5000:],
        )

    return process


def verify_output(
    output: Path,
    description: str,
):

    if not output.exists():

        raise HTTPException(
            status_code=422,
            detail=(
                f"{description} did not "
                "produce an output file."
            ),
        )

    if output.stat().st_size <= 0:

        raise HTTPException(
            status_code=422,
            detail=(
                f"{description} produced "
                "an empty output file."
            ),
        )

    return output


def file_response(
    output: Path,
    media_type: str,
    work: Path,
    engine: str,
):
    verify_output(
        output,
        "Conversion",
    )

    filename = output.name

    encoded_filename = quote(
        filename,
        safe="",
    )

    return FileResponse(
        path=str(output),
        media_type=media_type,
        filename=filename,
        headers={
            "Content-Disposition": (
                f'attachment; filename="{filename}"; '
                f"filename*=UTF-8''{encoded_filename}"
            ),
            "X-Converted-Filename": filename,
            "X-Conversion-Engine": engine,
            "Cache-Control": "no-store",
        },
        background=BackgroundTask(
            cleanup,
            work,
        ),
    )


# ============================================================
# OCR
# ============================================================


def normalize_ocr_language(
    language: str | None,
) -> str:

    value = (
        language
        or "eng"
    ).strip().lower()

    value = value.replace(
        "-",
        "_",
    )

    languages = {
        "eng": "eng",
        "english": "eng",

        "tam": "tam",
        "tamil": "tam",

        "hin": "hin",
        "hindi": "hin",

        "mal": "mal",
        "malayalam": "mal",

        "tel": "tel",
        "telugu": "tel",

        "kan": "kan",
        "kannada": "kan",

        "fra": "fra",
        "french": "fra",

        "deu": "deu",
        "ger": "deu",
        "german": "deu",

        "spa": "spa",
        "spanish": "spa",

        "ita": "ita",
        "italian": "ita",

        "por": "por",
        "portuguese": "por",
    }

    return languages.get(
        value,
        "eng",
    )


def tesseract_available() -> bool:

    return (
        shutil.which(
            "tesseract"
        )
        is not None
    )


def installed_tesseract_languages() -> set[str]:

    if not tesseract_available():
        return set()

    try:

        process = subprocess.run(
            [
                "tesseract",
                "--list-langs",
            ],
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            text=True,
            timeout=20,
            check=False,
        )

        lines = (
            process.stdout or ""
        ).splitlines()

        result = set()

        for line in lines:

            value = line.strip()

            if (
                value
                and not value.startswith(
                    "List of available languages"
                )
            ):
                result.add(
                    value
                )

        return result

    except Exception:

        return set()


def render_page_image(
    page,
    dpi: int = OCR_DPI,
):

    import fitz

    scale = dpi / 72.0

    matrix = fitz.Matrix(
        scale,
        scale,
    )

    pixmap = page.get_pixmap(
        matrix=matrix,
        alpha=False,
    )

    return pixmap, scale


def run_ocr(
    page,
    language: str = "eng",
):

    if not tesseract_available():

        raise HTTPException(
            status_code=503,
            detail=(
                "This PDF appears to be scanned, "
                "but Tesseract OCR is not installed "
                "on the conversion server."
            ),
        )

    try:

        import pytesseract
        from PIL import Image

    except ImportError as exc:

        raise HTTPException(
            status_code=503,
            detail=(
                "OCR dependencies are missing: "
                f"{exc}"
            ),
        )

    available = (
        installed_tesseract_languages()
    )

    requested_languages = (
        language.split("+")
    )

    usable_languages = [
        item
        for item in requested_languages
        if not available
        or item in available
    ]

    if not usable_languages:

        raise HTTPException(
            status_code=422,
            detail=(
                "The requested OCR language "
                f"'{language}' is not installed "
                "on the conversion server."
            ),
        )

    ocr_language = "+".join(
        usable_languages
    )

    pixmap, scale = (
        render_page_image(
            page
        )
    )

    image = Image.open(
        BytesIO(
            pixmap.tobytes(
                "png"
            )
        )
    )

    try:

        data = (
            pytesseract.image_to_data(
                image,
                lang=ocr_language,
                config="--oem 3 --psm 3",
                output_type=(
                    pytesseract.Output.DICT
                ),
            )
        )

    except Exception as exc:

        raise HTTPException(
            status_code=422,
            detail=(
                f"OCR failed: {exc}"
            ),
        )

    words = []

    count = len(
        data.get(
            "text",
            [],
        )
    )

    for index in range(count):

        text = (
            data["text"][index]
            or ""
        ).strip()

        if not text:
            continue

        try:

            confidence = float(
                data["conf"][index]
            )

        except (
            ValueError,
            TypeError,
        ):

            confidence = -1

        if confidence < 10:
            continue

        x = int(
            data["left"][index]
        )

        y = int(
            data["top"][index]
        )

        width = int(
            data["width"][index]
        )

        height = int(
            data["height"][index]
        )

        words.append(
            {
                "text": text,
                "x0": x / scale,
                "x1": (
                    x + width
                ) / scale,
                "top": y / scale,
                "bottom": (
                    y + height
                ) / scale,
                "height": max(
                    1.0,
                    height / scale,
                ),
                "confidence": confidence,
                "source": "ocr",
            }
        )

    return words


# ============================================================
# PDF NATIVE TEXT EXTRACTION
# ============================================================


def extract_pdf_words(
    page,
    language: str = "eng",
):

    words = []

    try:

        raw_words = (
            page.get_text(
                "words",
                sort=True,
            )
            or []
        )

    except Exception:

        raw_words = []

    for item in raw_words:

        if len(item) < 5:
            continue

        text = str(
            item[4]
        ).strip()

        if not text:
            continue

        x0 = float(
            item[0]
        )

        top = float(
            item[1]
        )

        x1 = float(
            item[2]
        )

        bottom = float(
            item[3]
        )

        words.append(
            {
                "text": text,
                "x0": x0,
                "x1": x1,
                "top": top,
                "bottom": bottom,
                "height": max(
                    1.0,
                    bottom - top,
                ),
                "source": "pdf",
            }
        )

    text_length = sum(
        len(
            word["text"]
        )
        for word in words
    )

    if (
        text_length
        < OCR_MIN_TEXT_CHARS
    ):

        try:

            ocr_words = run_ocr(
                page,
                language,
            )

            if ocr_words:
                return ocr_words

        except HTTPException:
            raise

        except Exception:
            pass

    return words


# ============================================================
# LINE RECONSTRUCTION
# ============================================================


def group_words_into_lines(
    words: list[dict[str, Any]],
    tolerance: float = 4.0,
):

    if not words:
        return []

    ordered = sorted(
        words,
        key=lambda word: (
            word["top"],
            word["x0"],
        ),
    )

    lines = []

    for word in ordered:

        center_y = (
            word["top"]
            + word["bottom"]
        ) / 2.0

        selected = None

        for line in reversed(
            lines[-8:]
        ):

            line_height = max(
                1.0,
                line["height"],
            )

            tolerance_value = max(
                tolerance,
                line_height * 0.55,
                word["height"] * 0.55,
            )

            if abs(
                center_y
                - line["center_y"]
            ) <= tolerance_value:

                selected = line
                break

        if selected is None:

            selected = {
                "center_y": center_y,
                "height": word[
                    "height"
                ],
                "words": [],
            }

            lines.append(
                selected
            )

        selected[
            "words"
        ].append(
            word
        )

        selected["height"] = max(
            selected["height"],
            word["height"],
        )

        selected["center_y"] = (
            selected["center_y"]
            + center_y
        ) / 2.0

    result = []

    for line in lines:

        line["words"].sort(
            key=lambda word:
            word["x0"]
        )

        if not line["words"]:
            continue

        line["text"] = " ".join(
            word["text"]
            for word in line[
                "words"
            ]
        ).strip()

        line["x0"] = min(
            word["x0"]
            for word in line[
                "words"
            ]
        )

        line["x1"] = max(
            word["x1"]
            for word in line[
                "words"
            ]
        )

        line["top"] = min(
            word["top"]
            for word in line[
                "words"
            ]
        )

        line["bottom"] = max(
            word["bottom"]
            for word in line[
                "words"
            ]
        )

        result.append(
            line
        )

    result.sort(
        key=lambda line: (
            line["top"],
            line["x0"],
        )
    )

    return result


def line_font_size(
    line: dict[str, Any],
) -> float:

    heights = [
        word["height"]
        for word in line[
            "words"
        ]
        if word.get("height")
    ]

    if not heights:
        return 10.0

    average = (
        sum(heights)
        / len(heights)
    )

    return max(
        6.0,
        min(
            32.0,
            average * 0.88,
        ),
    )


# ============================================================
# PDF → WORD
# ============================================================


def _pdf_font_name(name: str | None) -> str:
    """Map common PDF font names to fonts that are usually available in Word."""
    value = (name or "Arial").split("+")[-1].strip()
    lowered = value.lower()
    if "times" in lowered or "roman" in lowered:
        return "Times New Roman"
    if "courier" in lowered or "mono" in lowered:
        return "Courier New"
    if "arial" in lowered or "helvetica" in lowered or "calibri" in lowered:
        return "Arial"
    if "dejavu" in lowered:
        return "DejaVu Sans"
    return "Arial"


def _pdf_color_rgb(value: int | None):
    from docx.shared import RGBColor

    packed = int(value or 0)
    return RGBColor(
        (packed >> 16) & 0xFF,
        (packed >> 8) & 0xFF,
        packed & 0xFF,
    )


def _configure_docx_section(section, page):
    from docx.shared import Inches

    width_in = max(1.0, float(page.rect.width) / 72.0)
    height_in = max(1.0, float(page.rect.height) / 72.0)

    section.page_width = Inches(width_in)
    section.page_height = Inches(height_in)
    section.top_margin = Inches(0.25)
    section.bottom_margin = Inches(0.25)
    section.left_margin = Inches(0.25)
    section.right_margin = Inches(0.25)
    section.header_distance = Inches(0.1)
    section.footer_distance = Inches(0.1)


def _add_docx_text_block(document, block, page_width_pt: float):
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt
    from docx.oxml.ns import qn

    bbox = block.get("bbox") or (0, 0, page_width_pt, 0)
    x0, y0, x1, y1 = [float(v) for v in bbox]
    paragraph = document.add_paragraph()
    pf = paragraph.paragraph_format

    # Approximate the original PDF horizontal position.
    pf.left_indent = Inches(max(0.0, min((x0 / 72.0) - 0.25, max(0.0, page_width_pt / 72.0 - 0.5))))
    pf.space_before = Pt(0)
    pf.space_after = Pt(0)
    pf.line_spacing = 1.0

    lines = block.get("lines") or []
    for line_index, line in enumerate(lines):
        if line_index:
            paragraph.add_run().add_break()

        spans = line.get("spans") or []
        for span in spans:
            value = str(span.get("text") or "")
            if not value:
                continue

            run = paragraph.add_run(value)
            run.font.size = Pt(max(6.0, min(72.0, float(span.get("size") or 10.0))))
            font_name = _pdf_font_name(span.get("font"))
            run.font.name = font_name
            try:
                run._element.rPr.rFonts.set(qn("w:eastAsia"), font_name)
            except Exception:
                pass

            flags = int(span.get("flags") or 0)
            run.bold = bool(flags & 16)
            run.italic = bool(flags & 2)

            try:
                run.font.color.rgb = _pdf_color_rgb(span.get("color"))
            except Exception:
                pass

    # A paragraph containing only whitespace is not useful in DOCX.
    if not paragraph.text.strip():
        document._body._body.remove(paragraph._p)
        return False

    return True


def _add_docx_image_block(document, block, page_width_pt: float):
    from docx.shared import Inches

    image_bytes = block.get("image")
    if not image_bytes:
        return False

    bbox = block.get("bbox") or (0, 0, page_width_pt, page_width_pt)
    x0, y0, x1, y1 = [float(v) for v in bbox]
    width_pt = max(1.0, x1 - x0)
    max_width_in = max(1.0, page_width_pt / 72.0 - 0.5)
    width_in = min(width_pt / 72.0, max_width_in)

    paragraph = document.add_paragraph()
    paragraph.paragraph_format.left_indent = Inches(max(0.0, (x0 / 72.0) - 0.25))
    paragraph.paragraph_format.space_before = Inches(max(0.0, y0 / 72.0) * 0.0)
    paragraph.paragraph_format.space_after = Inches(0)

    try:
        run = paragraph.add_run()
        run.add_picture(BytesIO(image_bytes), width=Inches(width_in))
        return True
    except Exception:
        document._body._body.remove(paragraph._p)
        return False


def pdf_to_docx(
    source: Path,
    output: Path,
    language: str = "eng",
):
    """Convert PDF pages into a real DOCX with editable native text and embedded images.

    Native PDF text is reconstructed from PyMuPDF blocks/spans so font size, font family,
    bold/italic state, color and approximate horizontal placement survive. Scanned pages
    fall back to the existing Tesseract word reconstruction.
    """
    from docx import Document
    from docx.enum.section import WD_SECTION
    from docx.shared import Inches, Pt
    import fitz

    try:
        pdf = fitz.open(str(source))
    except Exception as exc:
        raise HTTPException(status_code=422, detail=f"Could not open PDF: {exc}")

    try:
        if pdf.page_count == 0:
            raise HTTPException(status_code=422, detail="The PDF has no pages.")
        if pdf.page_count > MAX_PDF_PAGES:
            raise HTTPException(
                status_code=413,
                detail=f"PDF contains more than {MAX_PDF_PAGES} pages.",
            )

        document = Document()
        total_content = 0

        for page_index, page in enumerate(pdf):
            if page_index == 0:
                section = document.sections[0]
            else:
                section = document.add_section(WD_SECTION.NEW_PAGE)
            _configure_docx_section(section, page)

            page_dict = page.get_text("dict", sort=True) or {}
            blocks = page_dict.get("blocks") or []
            native_text = sum(
                len(str(span.get("text") or ""))
                for block in blocks
                if block.get("type") == 0
                for line in (block.get("lines") or [])
                for span in (line.get("spans") or [])
            )

            if native_text >= OCR_MIN_TEXT_CHARS:
                for block in blocks:
                    block_type = block.get("type")
                    if block_type == 0:
                        if _add_docx_text_block(document, block, float(page.rect.width)):
                            total_content += 1
                    elif block_type == 1:
                        if _add_docx_image_block(document, block, float(page.rect.width)):
                            total_content += 1
            else:
                # Scanned/image-only page: use OCR words so the DOCX remains editable.
                words = extract_pdf_words(page, language)
                lines = group_words_into_lines(words)
                previous_bottom = None

                for line in lines:
                    text = str(line.get("text") or "").strip()
                    if not text:
                        continue

                    paragraph = document.add_paragraph()
                    pf = paragraph.paragraph_format
                    left = max(0.0, (float(line.get("x0", 0)) / 72.0) - 0.25)
                    pf.left_indent = Inches(min(left, max(0.0, page.rect.width / 72.0 - 0.5)))
                    if previous_bottom is not None:
                        gap = max(0.0, float(line.get("top", 0)) - previous_bottom)
                        pf.space_before = Pt(min(30.0, gap))
                    pf.space_after = Pt(0)
                    pf.line_spacing = 1.0

                    font_size = line_font_size(line)
                    for index, word in enumerate(line.get("words") or []):
                        if index:
                            paragraph.add_run(" ")
                        run = paragraph.add_run(str(word.get("text") or ""))
                        run.font.size = Pt(font_size)
                    previous_bottom = float(line.get("bottom", 0))
                    total_content += 1

            # Preserve an entirely blank page.
            if not blocks and native_text == 0:
                try:
                    words = extract_pdf_words(page, language)
                except Exception:
                    words = []
                if not words:
                    document.add_paragraph("")

        if total_content == 0:
            raise HTTPException(
                status_code=422,
                detail="No readable content was found in this PDF. If it is scanned, make sure the required Tesseract language is installed.",
            )

        output.parent.mkdir(parents=True, exist_ok=True)
        document.save(str(output))
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Could not create Word document: {exc}")
    finally:
        pdf.close()

    return verify_output(output, "PDF → Word")

def _pptx_font_name(name: str | None) -> str:
    value = (name or "Arial").split("+")[-1].strip()
    lowered = value.lower()
    if "times" in lowered or "roman" in lowered:
        return "Times New Roman"
    if "courier" in lowered or "mono" in lowered:
        return "Courier New"
    if "arial" in lowered or "helvetica" in lowered or "calibri" in lowered:
        return "Arial"
    if "dejavu" in lowered:
        return "DejaVu Sans"
    return "Arial"


def _pptx_color_rgb(value: int | None):
    from pptx.dml.color import RGBColor

    packed = int(value or 0)
    return RGBColor(
        (packed >> 16) & 0xFF,
        (packed >> 8) & 0xFF,
        packed & 0xFF,
    )


def _add_pptx_text_block(slide, block, page_rect, slide_width, slide_height):
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
    from pptx.util import Inches, Pt

    bbox = block.get("bbox")
    if not bbox:
        return False

    page_w = max(1.0, float(page_rect.width))
    page_h = max(1.0, float(page_rect.height))
    sx = float(slide_width) / page_w
    sy = float(slide_height) / page_h

    x0, y0, x1, y1 = [float(v) for v in bbox]
    left = max(0, int(x0 * sx))
    top = max(0, int(y0 * sy))
    width = max(1, int((x1 - x0) * sx))
    height = max(1, int((y1 - y0) * sy))

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    lines = block.get("lines") or []
    first_paragraph = True
    added = False

    for line in lines:
        paragraph = tf.paragraphs[0] if first_paragraph else tf.add_paragraph()
        first_paragraph = False
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 1.0

        spans = line.get("spans") or []
        for span in spans:
            value = str(span.get("text") or "")
            if not value:
                continue
            run = paragraph.add_run()
            run.text = value
            run.font.name = _pptx_font_name(span.get("font"))
            run.font.size = Pt(max(6.0, min(96.0, float(span.get("size") or 10.0))))
            flags = int(span.get("flags") or 0)
            run.font.bold = bool(flags & 16)
            run.font.italic = bool(flags & 2)
            try:
                run.font.color.rgb = _pptx_color_rgb(span.get("color"))
            except Exception:
                pass
            added = True

    if not added:
        slide.shapes._spTree.remove(shape._element)
        return False

    # PDF text is generally left aligned; preserve obvious centered/right blocks.
    block_center = (x0 + x1) / 2.0
    page_center = page_w / 2.0
    if abs(block_center - page_center) <= page_w * 0.08:
        for paragraph in tf.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER

    return True


def _add_pptx_image_block(slide, block, page_rect, slide_width, slide_height):
    from pptx.util import Inches

    image_bytes = block.get("image")
    bbox = block.get("bbox")
    if not image_bytes or not bbox:
        return False

    page_w = max(1.0, float(page_rect.width))
    page_h = max(1.0, float(page_rect.height))
    sx = float(slide_width) / page_w
    sy = float(slide_height) / page_h
    x0, y0, x1, y1 = [float(v) for v in bbox]

    left = int(x0 * sx)
    top = int(y0 * sy)
    width = max(1, int((x1 - x0) * sx))
    height = max(1, int((y1 - y0) * sy))

    try:
        slide.shapes.add_picture(BytesIO(image_bytes), left, top, width=width, height=height)
        return True
    except Exception:
        return False


def _add_pptx_ocr_page(slide, page, language, slide_width, slide_height):
    """Add editable OCR text to a slide for scanned PDFs."""
    words = extract_pdf_words(page, language)
    lines = group_words_into_lines(words)
    if not lines:
        return False

    page_rect = page.rect
    page_w = max(1.0, float(page_rect.width))
    page_h = max(1.0, float(page_rect.height))
    sx = float(slide_width) / page_w
    sy = float(slide_height) / page_h
    added = False

    from pptx.util import Pt

    for line in lines:
        x0 = float(line.get("x0", 0))
        y0 = float(line.get("top", 0))
        x1 = float(line.get("x1", x0 + 10))
        y1 = float(line.get("bottom", y0 + 12))
        shape = slide.shapes.add_textbox(
            int(x0 * sx),
            int(y0 * sy),
            max(1, int((x1 - x0) * sx)),
            max(1, int((y1 - y0) * sy)),
        )
        tf = shape.text_frame
        tf.clear()
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        run = tf.paragraphs[0].add_run()
        run.text = str(line.get("text") or "")
        run.font.name = "Arial"
        run.font.size = Pt(line_font_size(line))
        added = True

    return added


def pdf_to_pptx(
    source: Path,
    output: Path,
    language: str = "eng",
):
    """Convert PDF to an editable PowerPoint.

    Native PDF text becomes PowerPoint text boxes with editable runs. Native PDF raster
    images become PowerPoint pictures. Scanned pages use OCR text boxes; a page screenshot
    is used only when a page contains no extractable content at all.
    """
    from pptx import Presentation
    from pptx.util import Inches
    import fitz

    pdf = _fitz_open(source)
    try:
        if pdf.page_count == 0:
            raise HTTPException(status_code=422, detail="The PDF has no pages.")
        if pdf.page_count > MAX_PDF_PAGES:
            raise HTTPException(
                status_code=413,
                detail=f"PDF contains more than {MAX_PDF_PAGES} pages.",
            )

        first_rect = pdf[0].rect
        page_ratio = float(first_rect.width) / max(1.0, float(first_rect.height))

        # Keep the longest slide dimension at 13.333 in so portrait PDFs do not create
        # invalid PowerPoint dimensions.
        if page_ratio >= 1:
            slide_width_in = 13.333
            slide_height_in = 13.333 / page_ratio
        else:
            slide_height_in = 13.333
            slide_width_in = 13.333 * page_ratio

        presentation = Presentation()
        presentation.slide_width = Inches(slide_width_in)
        presentation.slide_height = Inches(slide_height_in)
        blank_layout = presentation.slide_layouts[6]

        for page in pdf:
            slide = presentation.slides.add_slide(blank_layout)
            page_dict = page.get_text("dict", sort=True) or {}
            blocks = page_dict.get("blocks") or []

            native_text = sum(
                len(str(span.get("text") or ""))
                for block in blocks
                if block.get("type") == 0
                for line in (block.get("lines") or [])
                for span in (line.get("spans") or [])
            )

            added = False
            if native_text >= OCR_MIN_TEXT_CHARS:
                for block in blocks:
                    if block.get("type") == 0:
                        added = _add_pptx_text_block(
                            slide,
                            block,
                            page.rect,
                            presentation.slide_width,
                            presentation.slide_height,
                        ) or added
                    elif block.get("type") == 1:
                        added = _add_pptx_image_block(
                            slide,
                            block,
                            page.rect,
                            presentation.slide_width,
                            presentation.slide_height,
                        ) or added
            else:
                added = _add_pptx_ocr_page(
                    slide,
                    page,
                    language,
                    presentation.slide_width,
                    presentation.slide_height,
                )

            # If a PDF page truly contains no readable text/image blocks, retain it as
            # a visual page rather than returning an empty slide.
            if not added:
                pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                slide.shapes.add_picture(
                    BytesIO(pixmap.tobytes("png")),
                    0,
                    0,
                    width=presentation.slide_width,
                    height=presentation.slide_height,
                )

        output.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(output))
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Could not create PowerPoint: {exc}")
    finally:
        pdf.close()

    return verify_output(output, "PDF → PowerPoint")

def _fitz_open(
    source: Path,
):

    try:

        import fitz

        return fitz.open(
            str(source)
        )

    except Exception as exc:

        raise HTTPException(
            status_code=422,
            detail=(
                f"Could not read PDF: "
                f"{exc}"
            ),
        )


# ============================================================
# OFFICE → PDF
# ============================================================


def office_to_pdf(
    source: Path,
    outdir: Path,
    profile: Path,
):

    binary = binary_path(
        "soffice",
        "libreoffice",
    )

    outdir.mkdir(
        parents=True,
        exist_ok=True,
    )

    profile.mkdir(
        parents=True,
        exist_ok=True,
    )

    # LibreOffice requires a unique user profile
    # for concurrent headless conversions.
    profile_uri = (
        "file://"
        + profile.resolve().as_posix()
    )

    args = [
        binary,
        "--headless",
        "--invisible",
        "--nodefault",
        "--nofirststartwizard",
        "--nologo",
        "--convert-to",
        "pdf",
        "--outdir",
        str(outdir),
        f"-env:UserInstallation={profile_uri}",
        str(source),
    ]

    run_checked(
        args,
        timeout=240,
    )

    output = (
        outdir
        / f"{source.stem}.pdf"
    )

    return verify_output(
        output,
        "Office → PDF",
    )


# ============================================================
# HTML → PDF
# ============================================================


def html_to_pdf(
    source: Path,
    outdir: Path,
    profile: Path,
):

    # LibreOffice is used as the deterministic
    # server-side renderer.
    #
    # For HTML requiring full Chromium CSS/JS,
    # Playwright can be added separately.

    return office_to_pdf(
        source,
        outdir,
        profile,
    )


# ============================================================
# QPDF
# ============================================================


def qpdf_transform(
    source: Path,
    output: Path,
    password: str | None,
    mode: str,
):

    qpdf = binary_path(
        "qpdf"
    )

    output.parent.mkdir(
        parents=True,
        exist_ok=True,
    )

    if mode == "unlock":

        args = [
            qpdf,
            "--password="
            + (
                password
                or ""
            ),
            "--decrypt",
            str(source),
            str(output),
        ]

    elif mode == "protect":

        if not password:

            raise HTTPException(
                status_code=400,
                detail=(
                    "A password is required "
                    "to protect the PDF."
                ),
            )

        args = [
            qpdf,
            "--encrypt",
            password,
            password,
            "256",
            "--",
            str(source),
            str(output),
        ]

    else:

        raise HTTPException(
            status_code=400,
            detail=(
                "Unsupported qpdf operation."
            ),
        )

    run_checked(
        args,
        timeout=240,
    )

    return verify_output(
        output,
        "PDF security operation",
    )


# ============================================================
# PDF/A
# ============================================================


def find_ghostscript_icc() -> str | None:

    candidates = [
        os.getenv(
            "GS_ICC_PROFILE"
        ),

        "/usr/share/color/icc/ghostscript/srgb.icc",

        "/usr/share/ghostscript/iccprofiles/srgb.icc",
    ]

    if os.name == "nt":

        candidates.extend(
            [
                r"C:\Program Files\gs\gs10.00.0\iccprofiles\srgb.icc",
                r"C:\Program Files\gs\gs10.01.0\iccprofiles\srgb.icc",
                r"C:\Program Files\gs\gs10.02.0\iccprofiles\srgb.icc",
                r"C:\Program Files\gs\gs10.03.0\iccprofiles\srgb.icc",
                r"C:\Program Files\gs\gs10.04.0\iccprofiles\srgb.icc",
                r"C:\Program Files\gs\gs10.05.0\iccprofiles\srgb.icc",
            ]
        )

    for candidate in candidates:

        if candidate and Path(
            candidate
        ).exists():

            return candidate

    return None


def pdf_to_pdfa(
    source: Path,
    output: Path,
):

    gs = binary_path(
        "gs",
        "gswin64c",
        "gswin32c",
    )

    icc_profile = (
        find_ghostscript_icc()
    )

    if not icc_profile:

        raise HTTPException(
            status_code=503,
            detail=(
                "Ghostscript is installed, but "
                "an sRGB ICC profile could not be "
                "found for PDF/A conversion."
            ),
        )

    output.parent.mkdir(
        parents=True,
        exist_ok=True,
    )

    args = [
        gs,
        "-dPDFA=2",
        "-dBATCH",
        "-dNOPAUSE",
        "-dSAFER",
        "-sDEVICE=pdfwrite",
        "-sColorConversionStrategy=RGB",
        "-sProcessColorModel=DeviceRGB",
        f"-sOutputICCProfile={icc_profile}",
        "-dPDFACompatibilityPolicy=1",
        "-dAutoRotatePages=/None",
        "-o",
        str(output),
        str(source),
    ]

    run_checked(
        args,
        timeout=300,
    )

    return verify_output(
        output,
        "PDF/A conversion",
    )


# ============================================================
# PDF TEXT
# ============================================================


def extract_pdf_text(
    source: Path,
) -> str:

    import fitz

    pdf = _fitz_open(
        source
    )

    try:

        pages = []

        for page in pdf:

            text = (
                page.get_text(
                    "text"
                )
                or ""
            )

            pages.append(
                text.strip()
            )

        return "\n\n".join(
            page
            for page in pages
            if page
        )

    finally:

        pdf.close()


# ============================================================
# TRANSLATION
# ============================================================


def translation_request(
    text: str,
    target_lang: str,
) -> str:

    endpoint = os.getenv(
        "TRANSLATION_API_URL"
    )

    if not endpoint:

        raise HTTPException(
            status_code=503,
            detail=(
                "TRANSLATION_API_URL is not "
                "configured on the conversion server."
            ),
        )

    payload = json.dumps(
        {
            "q": text,
            "target": target_lang,
            "source": "auto",
            "format": "text",
        }
    ).encode(
        "utf-8"
    )

    request = urllib.request.Request(
        endpoint,
        data=payload,
        headers={
            "Content-Type":
                "application/json",
            "Accept":
                "application/json",
        },
        method="POST",
    )

    try:

        with urllib.request.urlopen(
            request,
            timeout=120,
        ) as response:

            body = json.loads(
                response.read().decode(
                    "utf-8"
                )
            )

    except Exception as exc:

        raise HTTPException(
            status_code=502,
            detail=(
                f"Translation service failed: "
                f"{exc}"
            ),
        )

    translated = (
        body.get(
            "translatedText"
        )
        or body.get(
            "translation"
        )
        or body.get(
            "text"
        )
        or ""
    )

    return str(
        translated
    )


# ============================================================
# SEND EMAIL
# ============================================================


@app.post(
    "/send-email"
)
async def send_email(
    file: UploadFile = File(...),
    to: str = Form(...),
    subject: str = Form(...),
    tool: str = Form(...),
):

    api_key = os.getenv(
        "RESEND_API_KEY"
    )

    email_from = os.getenv(
        "EMAIL_FROM",
        "QuadraConverter <onboarding@resend.dev>",
    )

    if not api_key:

        raise HTTPException(
            status_code=503,
            detail=(
                "Email service is not configured."
            ),
        )

    if (
        not to
        or "@"
        not in to
    ):

        raise HTTPException(
            status_code=400,
            detail=(
                "Please provide a valid email address."
            ),
        )

    content = await file.read()

    if not content:

        raise HTTPException(
            status_code=400,
            detail=(
                "The converted file is empty."
            ),
        )

    if (
        len(content)
        > MAX_FILE_BYTES
    ):

        raise HTTPException(
            status_code=413,
            detail=(
                "The file is too large to email."
            ),
        )

    try:

        import resend

        resend.api_key = (
            api_key
        )

        attachment = {
            "filename": (
                file.filename
                or "converted-file"
            ),
            "content": list(
                content
            ),
        }

        response = (
            resend.Emails.send(
                {
                    "from": email_from,
                    "to": [to],
                    "subject": subject,
                    "html": f"""
                        <div style="font-family:Arial,sans-serif">
                            <h2>QuadraConverter</h2>

                            <p>
                                Your file has been converted
                                successfully using
                                <strong>{tool}</strong>.
                            </p>

                            <p>
                                The converted file is attached
                                to this email.
                            </p>

                            <p>
                                — QuadraConverter
                            </p>
                        </div>
                    """,
                    "attachments": [
                        attachment
                    ],
                }
            )
        )

        return JSONResponse(
            {
                "success": True,
                "message": (
                    "Email sent successfully."
                ),
                "id": (
                    response.get("id")
                    if isinstance(
                        response,
                        dict,
                    )
                    else None
                ),
            }
        )

    except HTTPException:
        raise

    except Exception as exc:

        raise HTTPException(
            status_code=502,
            detail=(
                f"Email provider failed: "
                f"{exc}"
            ),
        )


# ============================================================
# HEALTH CHECK
# ============================================================


@app.get(
    "/health"
)
def health():

    engines = []

    checks = [
        (
            "LibreOffice",
            (
                "soffice",
                "libreoffice",
            ),
        ),
        (
            "qpdf",
            (
                "qpdf",
            ),
        ),
        (
            "Ghostscript",
            (
                "gs",
                "gswin64c",
                "gswin32c",
            ),
        ),
        (
            "Tesseract",
            (
                "tesseract",
            ),
        ),
    ]

    for label, names in checks:

        try:

            path = binary_path(
                *names
            )

            engines.append(
                {
                    "name": label,
                    "available": True,
                    "binary": path,
                }
            )

        except Exception:

            engines.append(
                {
                    "name": label,
                    "available": False,
                    "binary": None,
                }
            )

    return {
        "ok": True,
        "service": APP_NAME,
        "version": APP_VERSION,
        "max_file_mb": (
            MAX_FILE_BYTES
            // (1024 * 1024)
        ),
        "ocr_dpi": OCR_DPI,
        "engines": engines,
    }


# ============================================================
# SINGLE CONVERSION ENDPOINT
#
# DO NOT CREATE ANOTHER @app.post("/convert") BELOW THIS.
# ============================================================


@app.post(
    "/convert"
)
async def convert(
    file: UploadFile = File(...),
    operation: str = Form(...),
    password: str = Form(""),
    targetLang: str = Form(""),
    language: str = Form("eng"),
):

    work = Path(
        tempfile.mkdtemp(
            prefix="quadra-convert-"
        )
    )

    try:

        if not file.filename:

            raise HTTPException(
                status_code=400,
                detail=(
                    "No file was uploaded."
                ),
            )

        operation = (
            operation
            or ""
        ).strip().lower()

        language = (
            language
            or "eng"
        ).strip()

        # ====================================================
        # OFFICE → PDF
        # ====================================================

        if operation == "office-to-pdf":

            source = save_upload(
                file,
                work,
                ALLOWED_OFFICE,
            )

            outdir = (
                work / "out"
            )

            profile = (
                work / "profile"
            )

            outdir.mkdir(
                parents=True,
                exist_ok=True,
            )

            profile.mkdir(
                parents=True,
                exist_ok=True,
            )

            output = office_to_pdf(
                source,
                outdir,
                profile,
            )

            return file_response(
                output,
                "application/pdf",
                work,
                "LibreOffice",
            )

        # ====================================================
        # HTML → PDF
        # ====================================================

        if operation == "html-to-pdf":

            source = save_upload(
                file,
                work,
                ALLOWED_HTML,
            )

            outdir = (
                work / "out"
            )

            profile = (
                work / "profile"
            )

            outdir.mkdir(
                parents=True,
                exist_ok=True,
            )

            profile.mkdir(
                parents=True,
                exist_ok=True,
            )

            output = html_to_pdf(
                source,
                outdir,
                profile,
            )

            return file_response(
                output,
                "application/pdf",
                work,
                "LibreOffice",
            )

        # ====================================================
        # PDF UNLOCK
        # ====================================================

        if operation == "pdf-unlock":

            source = save_upload(
                file,
                work,
                ALLOWED_PDF,
            )

            outdir = (
                work / "out"
            )

            outdir.mkdir(
                parents=True,
                exist_ok=True,
            )

            output = qpdf_transform(
                source,
                (
                    outdir
                    / f"{source.stem}-unlocked.pdf"
                ),
                password,
                "unlock",
            )

            return file_response(
                output,
                "application/pdf",
                work,
                "qpdf",
            )

        # ====================================================
        # PDF PROTECT
        # ====================================================

        if operation == "pdf-protect":

            if not password:

                raise HTTPException(
                    status_code=400,
                    detail=(
                        "A password is required "
                        "to protect the PDF."
                    ),
                )

            source = save_upload(
                file,
                work,
                ALLOWED_PDF,
            )

            outdir = (
                work / "out"
            )

            outdir.mkdir(
                parents=True,
                exist_ok=True,
            )

            output = qpdf_transform(
                source,
                (
                    outdir
                    / f"{source.stem}-protected.pdf"
                ),
                password,
                "protect",
            )

            return file_response(
                output,
                "application/pdf",
                work,
                "qpdf",
            )

        # ====================================================
        # PDF → PDF/A
        # ====================================================

        if operation == "pdf-to-pdfa":

            source = save_upload(
                file,
                work,
                ALLOWED_PDF,
            )

            outdir = (
                work / "out"
            )

            outdir.mkdir(
                parents=True,
                exist_ok=True,
            )

            output = pdf_to_pdfa(
                source,
                (
                    outdir
                    / f"{source.stem}-pdfa.pdf"
                ),
            )

            return file_response(
                output,
                "application/pdf",
                work,
                "Ghostscript",
            )

        # ====================================================
        # PDF → WORD
        #
        # IMPORTANT:
        # This is NOT browser-side conversion.
        #
        # Pipeline:
        #
        # PDF
        #  ↓
        # PyMuPDF native text
        #  ↓
        # coordinate extraction
        #  ↓
        # line reconstruction
        #  ↓
        # scanned-page detection
        #  ↓
        # Tesseract OCR when required
        #  ↓
        # Word paragraph reconstruction
        #  ↓
        # DOCX
        # ====================================================

        if operation == "pdf-to-word":

            source = save_upload(
                file,
                work,
                ALLOWED_PDF,
            )

            outdir = (
                work / "out"
            )

            outdir.mkdir(
                parents=True,
                exist_ok=True,
            )

            ocr_language = (
                normalize_ocr_language(
                    language
                )
            )

            output = pdf_to_docx(
                source,
                (
                    outdir
                    / f"{source.stem}.docx"
                ),
                ocr_language,
            )

            return file_response(
                output,
                (
                    "application/"
                    "vnd.openxmlformats-officedocument."
                    "wordprocessingml.document"
                ),
                work,
                "PyMuPDF + Tesseract + python-docx",
            )

        # ====================================================
        # PDF → EXCEL
        #
        # Pipeline:
        #
        # PDF
        #  ↓
        # native table detection
        #  ↓
        # coordinate-based text extraction
        #  ↓
        # column clustering
        #  ↓
        # OCR fallback for scanned PDF
        #  ↓
        # XLSX
        # ====================================================

        if operation == "pdf-to-xlsx":

            source = save_upload(
                file,
                work,
                ALLOWED_PDF,
            )

            outdir = (
                work / "out"
            )

            outdir.mkdir(
                parents=True,
                exist_ok=True,
            )

            ocr_language = (
                normalize_ocr_language(
                    language
                )
            )

            output = pdf_to_xlsx(
                source,
                (
                    outdir
                    / f"{source.stem}.xlsx"
                ),
                ocr_language,
            )

            return file_response(
                output,
                (
                    "application/"
                    "vnd.openxmlformats-officedocument."
                    "spreadsheetml.sheet"
                ),
                work,
                "pdfplumber + PyMuPDF + Tesseract + openpyxl",
            )

        # ====================================================
        # PDF → POWERPOINT
        # ====================================================

        if operation == "pdf-to-pptx":

            source = save_upload(
                file,
                work,
                ALLOWED_PDF,
            )

            outdir = (
                work / "out"
            )

            outdir.mkdir(
                parents=True,
                exist_ok=True,
            )

            output = pdf_to_pptx(
                source,
                (
                    outdir
                    / f"{source.stem}.pptx"
                ),
                normalize_ocr_language(language),
            )

            return file_response(
                output,
                (
                    "application/"
                    "vnd.openxmlformats-officedocument."
                    "presentationml.presentation"
                ),
                work,
                "PyMuPDF + python-pptx",
            )

        # ====================================================
        # PDF TRANSLATE
        # ====================================================

        if operation == "pdf-translate":

            source = save_upload(
                file,
                work,
                ALLOWED_PDF,
            )

            text = extract_pdf_text(
                source
            )

            if not text.strip():

                raise HTTPException(
                    status_code=422,
                    detail=(
                        "No readable text was found "
                        "in this PDF."
                    ),
                )

            target_language = (
                targetLang
                or "en"
            ).strip()

            translated = (
                translation_request(
                    text,
                    target_language,
                )
            )

            if not translated.strip():

                raise HTTPException(
                    status_code=422,
                    detail=(
                        "The translation service "
                        "returned no translated text."
                    ),
                )

            output = (
                work
                / f"{source.stem}-translated.txt"
            )

            output.write_text(
                translated,
                encoding="utf-8",
            )

            return file_response(
                output,
                "text/plain; charset=utf-8",
                work,
                "Translation API",
            )

        # ====================================================
        # UNKNOWN OPERATION
        # ====================================================

        raise HTTPException(
            status_code=400,
            detail=(
                f"Unsupported conversion operation: "
                f"{operation}"
            ),
        )

    except HTTPException:

        cleanup(
            work
        )

        raise

    except Exception as exc:

        cleanup(
            work
        )

        print(
            "[QuadraConverter] "
            f"Unexpected conversion error: {exc}"
        )

        raise HTTPException(
            status_code=500,
            detail=(
                "Conversion failed: "
                f"{str(exc)}"
            ),
        )

    finally:

        try:
            await file.close()
        except Exception:
            pass


# ============================================================
# STARTUP DIAGNOSTICS
# ============================================================


@app.on_event(
    "startup"
)
def startup_check():

    print("")
    print(
        "============================================================"
    )
    print(
        " QuadraConverter Conversion API"
    )
    print(
        f" Version: {APP_VERSION}"
    )
    print(
        "============================================================"
    )

    engines = [
        (
            "LibreOffice",
            (
                "soffice",
                "libreoffice",
            ),
        ),
        (
            "qpdf",
            (
                "qpdf",
            ),
        ),
        (
            "Ghostscript",
            (
                "gs",
                "gswin64c",
                "gswin32c",
            ),
        ),
        (
            "Tesseract",
            (
                "tesseract",
            ),
        ),
    ]

    for label, names in engines:

        try:

            path = binary_path(
                *names
            )

            print(
                f"[OK] {label}: {path}"
            )

        except Exception:

            print(
                f"[WARN] {label}: NOT INSTALLED"
            )

    print(
        "============================================================"
    )
    print(
        f"Max file size: "
        f"{MAX_FILE_BYTES // (1024 * 1024)} MB"
    )
    print(
        f"OCR DPI: {OCR_DPI}"
    )
    print(
        f"OCR threshold: "
        f"{OCR_MIN_TEXT_CHARS} characters"
    )
    print(
        "============================================================"
    )
    print("")


# ============================================================
# LOCAL DEVELOPMENT
# ============================================================

if __name__ == "__main__":

    import uvicorn

    uvicorn.run(
        "converter_api:app",
        host="0.0.0.0",
        port=int(
            os.getenv(
                "PORT",
                "8000",
            )
        ),
        reload=False,
    )
